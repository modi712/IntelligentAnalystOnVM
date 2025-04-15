import os
import uuid
import datetime
import logging
from pathlib import Path
from django.conf import settings
from uuid import uuid4
from pptx import Presentation
from pdf2image import convert_from_path
import tempfile
import shutil
from langchain_community.document_loaders import PyPDFLoader, CSVLoader
from langchain_huggingface.embeddings import HuggingFaceEmbeddings
from langchain.embeddings import SentenceTransformerEmbeddings
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser, JsonOutputParser
from langchain_core.prompts import PromptTemplate
from langchain_chroma import Chroma
from langchain_core.pydantic_v1 import BaseModel, Field, validator
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains import LLMChain
import json



# Configure logger
logger = logging.getLogger(__name__)

# Initialize constants
MEDIA_ROOT = settings.MEDIA_ROOT
REPORTS_DIR = os.path.join(MEDIA_ROOT, "reports")
VECTOR_STORES_DIR = os.path.join(MEDIA_ROOT, "vector_stores")
PPT_TEMPLATE_PATH = os.path.join(settings.BASE_DIR, "chatbot", "templates", "ppt_templates", "report_template.pptx")

# Initialize embeddings and text splitter
embeddings = SentenceTransformerEmbeddings(model_name="all-MiniLM-L6-v2")
text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)

# Initialize LLM for report generation
try:
    from .api_key import GROQ_API_KEY  # Import API key
    from langchain_groq import ChatGroq as Groq
    llm = Groq(api_key='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', model_name="llama-3.3-70b-versatile")
except (ImportError, Exception) as e:
    logger.error(f"Failed to initialize LLM: {e}")
    llm = None

# Create directory structure functions
def ensure_dir_exists(path):
    """Ensure the directory exists, create if it doesn't."""
    if not os.path.exists(path):
        os.makedirs(path)
    return path

def get_company_dir(company_name):
    """Get company directory path and create if doesn't exist."""
    company_path = os.path.join(VECTOR_STORES_DIR, company_name)
    return ensure_dir_exists(company_path)

def get_kb_dir(company_name, kb_name):
    """Get knowledge base directory path and create if doesn't exist."""
    kb_path = os.path.join(get_company_dir(company_name), kb_name)
    return ensure_dir_exists(kb_path)

def get_ppt_dir(company_name, kb_name):
    """Get PowerPoint directory path and create if doesn't exist."""
    ppt_path = os.path.join(get_kb_dir(company_name, kb_name), "ppt")
    return ensure_dir_exists(ppt_path)

def get_pdf_dir(company_name, kb_name):
    """Get PDF directory path and create if doesn't exist."""
    pdf_path = os.path.join(get_kb_dir(company_name, kb_name), "pdf")
    return ensure_dir_exists(pdf_path)

# Add these helper functions to make it easier to target specific shapes in your template:
def update_shape_text(ppt, slide_index, shape_index, text):
    """Update the text in a specific shape on a slide"""
    try:
        if slide_index < len(ppt.slides):
            slide = ppt.slides[slide_index]
            if shape_index < len(slide.shapes):
                shape = slide.shapes[shape_index]
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.text = text
                    return True
        return False
    except Exception as e:
        logger.error(f"Error updating shape text: {e}")
        return False

def runquery(query_text, slide_index, shape_index, page_num, paragraph_index, retriever, ppt=None):
    """Run a query and update a specific shape in the presentation"""
    try:
        # Get the response from the query
        response = run_query(query_text, retriever)
        
        # If no presentation is provided, just return the response
        if ppt is None:
            return response
        
        # Get the slide
        if slide_index < len(ppt.slides):
            slide = ppt.slides[slide_index]
            
            # Update the shape
            if shape_index < len(slide.shapes):
                shape = slide.shapes[shape_index]
                
                if hasattr(shape, 'text_frame'):
                    # If paragraph_index is specified, update only that paragraph
                    if paragraph_index < len(shape.text_frame.paragraphs):
                        shape.text_frame.paragraphs[paragraph_index].text = response
                    else:
                        shape.text_frame.text = response
                    
                    return True
        return False
    except Exception as e:
        logger.error(f"Error in runquery: {e}")
        return False

def page4sec3(retriever, ppt):
    """Update summary metrics on page 4 section 3"""
    runquery("Summarize Bahrain's GDP growth over years in around 30 words.", 3, 0, 4, 0, retriever, ppt)
    runquery("Summarize Bahrain's Inflation rate over years in around 30 words.", 3, 0, 4, 1, retriever, ppt)
    runquery("Summarize Bahrain's Unemployment rate over years in around 30 words.", 3, 0, 4, 2, retriever, ppt)
    runquery("Summarize Bahrain's Fiscal Balance over years in around 30 words.", 3, 0, 4, 3, retriever, ppt)
    runquery("Summarize Bahrain's Debt to GDP ratio over years in around 30 words.", 3, 0, 4, 4, retriever, ppt)
    runquery("Summarize Bahrain's Net Issuance over years in around 30 words.", 3, 0, 4, 5, retriever, ppt)

def page4sec1(retriever, ppt):
    """Update fiscal indicators on page 3"""
    runquery("What is fiscal balance as percentage of GDP for bahrain over years, answer in one paragraph of 60 words.", 2, 0, 3, 0, retriever, ppt)
    runquery("What is the Debt to GDP Ratio for Bahrain? Answer in 60 words.", 2, 0, 3, 1, retriever, ppt)
    runquery("Analyze Net Issuance vs Fiscal Deficit for Bahrain in 60 words.", 2, 0, 3, 2, retriever, ppt)

def page3sec1(retriever, ppt):
    """Update economic indicators on page 2"""
    runquery("What is GDP growth of Bahrain over years? Return answer in single paragraph of around 60 words.", 1, 1, 2, 0, retriever, ppt)
    runquery("What is inflation rate of Bahrain over years? Return answer in single paragraph of around 60 words.", 1, 1, 2, 1, retriever, ppt)
    runquery("What is unemployment rate of Bahrain over years? Return answer in single paragraph of around 60 words.", 1, 1, 2, 2, retriever, ppt)

# Process PDF files and create vector store
def create_vector_store(files, kb_name, company_name):
    """Create a vector store from uploaded files."""
    logger.info(f"Creating vector store for {company_name}/{kb_name}")
    
    # Get directories
    kb_dir = get_kb_dir(company_name, kb_name)
    ppt_dir = get_ppt_dir(company_name, kb_name)
    
    # Save uploaded files to directory
    file_paths = []
    for file_obj in files:
        # Save file to disk
        file_path = os.path.join(kb_dir, file_obj.name)
        with open(file_path, "wb") as f:
            for chunk in file_obj.chunks():
                f.write(chunk)
        file_paths.append(file_path)
        logger.debug(f"Saved file: {file_path}")
    
    # Process PDF files
    all_texts = []
    for file_path in file_paths:
        if file_path.lower().endswith('.pdf'):
            try:
                loader = PyPDFLoader(file_path)
                pages = loader.load_and_split()
                all_texts.extend([page.page_content for page in pages])
                logger.debug(f"Processed PDF: {file_path}")
            except Exception as e:
                logger.error(f"Error processing PDF {file_path}: {e}")
    
    # Create chunks and vectorize
    chunks = text_splitter.create_documents(all_texts)
    uuids = [str(uuid4()) for _ in range(len(chunks))]
    
    # Initialize and save vector store
    persist_dir = os.path.join(kb_dir, "vector_store")
    ensure_dir_exists(persist_dir)
    
    vectorstore = Chroma(
        collection_name=kb_name,
        embedding_function=embeddings,
        persist_directory=persist_dir,
    )
    
    # Add documents to the vector store
    try:
        vectorstore.add_documents(documents=chunks, ids=uuids)
        logger.info(f"Added {len(chunks)} document chunks to vector store")
    except Exception as e:
        logger.error(f"Error adding documents to vector store: {e}")
    
    # Create retriever for similarity search
    retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 3})
    
    # Generate report if company is CBB
    # In the create_vector_store function, modify the CBB report generation:
    if company_name == "Central Bank of Bahrain":
        try:
            # Generate the CBB report using the template and specific functions
            generate_cbb_report(retriever, company_name, kb_name)
            
            # Call the specific page functions for more detailed content
            page3sec1(retriever)  # Economic indicators
            page4sec1(retriever)  # Fiscal indicators
            page4sec3(retriever)  # Summary metrics
        except Exception as e:
            logger.error(f"Error generating CBB report: {e}")
    else:
        generate_generic_report(retriever, company_name, kb_name)

def generate_cbb_report(retriever, company_name, kb_name):
    """Generate detailed report for Central Bank of Bahrain with specific sections"""
    logger.info(f"Generating CBB report for {company_name}/{kb_name}")
    
    # Load the custom template
    try:
        ppt = Presentation(PPT_TEMPLATE_PATH)
    except Exception as e:
        logger.error(f"Error loading template: {e}")
        # Fall back to creating a new presentation
        ppt = Presentation()
    
    # Generate specific content for each section
    try:
        # Page 1: Strengths and Weaknesses
        strengths_weaknesses = run_query(
            "What are three Macroeconomic strengths and three weaknesses for Bahrain? Format as bullet points with 15 words for each point.",
            retriever
        )
        # Add to Page 1, Shape 3
        if len(ppt.slides) >= 1:
            slide = ppt.slides[0]
            if len(slide.shapes) >= 3:
                shape = slide.shapes[2]
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.text = strengths_weaknesses
        
        # Page 2: Macroeconomic Overview
        macro_overview = run_query(
            "Provide a comprehensive 250-word overview of Bahrain's economy.",
            retriever
        )
        # Add to Page 2, Shape 1
        if len(ppt.slides) >= 2:
            slide = ppt.slides[1]
            if len(slide.shapes) >= 1:
                shape = slide.shapes[0]
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.text = macro_overview
        
        # Page 2: Key Economic Indicators
        gdp_growth_analysis = run_query(
            "Analyze Bahrain's GDP Growth in 60 words.",
            retriever
        )
        inflation_analysis = run_query(
            "Analyze Bahrain's Inflation Rate in 60 words.",
            retriever
        )
        unemployment_analysis = run_query(
            "Analyze Bahrain's Unemployment Rate in 60 words.",
            retriever
        )
        
        # Combine economic indicators and add to Page 2, Shape 1
        economic_indicators = f"GDP Growth:\n{gdp_growth_analysis}\n\nInflation Rate:\n{inflation_analysis}\n\nUnemployment Rate:\n{unemployment_analysis}"
        if len(ppt.slides) >= 2:
            slide = ppt.slides[1]
            if len(slide.shapes) >= 2:
                shape = slide.shapes[1]
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.text = economic_indicators
        
        # Page 3: Fiscal Indicators
        fiscal_balance = run_query(
            "What is fiscal balance as percentage of GDP for Bahrain over years? Return answer in single paragraph of around 60 words.",
            retriever
        )
        debt_gdp = run_query(
            "What is the Debt to GDP Ratio for Bahrain? Answer in 60 words.",
            retriever
        )
        net_issuance = run_query(
            "Analyze Net Issuance vs Fiscal Deficit for Bahrain in 60 words.",
            retriever
        )
        
        # Add fiscal indicators to Page 3
        fiscal_indicators = f"Fiscal Balance:\n{fiscal_balance}\n\nDebt to GDP Ratio:\n{debt_gdp}\n\nNet Issuance vs Fiscal Deficit:\n{net_issuance}"
        if len(ppt.slides) >= 3:
            slide = ppt.slides[2]
            if len(slide.shapes) >= 1:
                shape = slide.shapes[0]
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.text = fiscal_indicators
        
        # Page 3/4: Summary Metrics (30 words each)
        gdp_summary = run_query("Summarize Bahrain's GDP growth over years in around 30 words.", retriever)
        inflation_summary = run_query("Summarize Bahrain's Inflation rate over years in around 30 words.", retriever)
        unemployment_summary = run_query("Summarize Bahrain's Unemployment rate over years in around 30 words.", retriever)
        fiscal_balance_summary = run_query("Summarize Bahrain's Fiscal Balance over years in around 30 words.", retriever)
        debt_gdp_summary = run_query("Summarize Bahrain's Debt to GDP ratio over years in around 30 words.", retriever)
        net_issuance_summary = run_query("Summarize Bahrain's Net Issuance over years in around 30 words.", retriever)
        
        # Add summary metrics to the appropriate slide/shape (assuming Page 3, Shape 1)
        summary_metrics = f"""GDP Growth: {gdp_summary}
        
Inflation Rate: {inflation_summary}

Unemployment Rate: {unemployment_summary}

Fiscal Balance: {fiscal_balance_summary}

Debt to GDP Ratio: {debt_gdp_summary}

Net Issuance: {net_issuance_summary}"""
        
        if len(ppt.slides) >= 3:
            slide = ppt.slides[2]
            if len(slide.shapes) >= 2:
                shape = slide.shapes[1]
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.text = summary_metrics
        
        # Save the presentation
        save_presentation(ppt, company_name, kb_name)
        
    except Exception as e:
        logger.error(f"Error generating CBB report: {e}")

def create_vector_store1(files, kb_name, company_name):
    """Create a vector store from uploaded files."""
    logger.info(f"Creating vector store for {company_name}/{kb_name}")
    
    # Get directories
    kb_dir = get_kb_dir(company_name, kb_name)
    ppt_dir = get_ppt_dir(company_name, kb_name)
    
    # Save uploaded files to directory
    file_paths = []
    for file_obj in files:
        # Save file to disk
        file_path = os.path.join(kb_dir, file_obj.name)
        with open(file_path, "wb") as f:
            for chunk in file_obj.chunks():
                f.write(chunk)
        file_paths.append(file_path)
        logger.debug(f"Saved file: {file_path}")
    
    # Process PDF files
    all_texts = []
    for file_path in file_paths:
        if file_path.lower().endswith('.pdf'):
            try:
                loader = PyPDFLoader(file_path)
                pages = loader.load_and_split()
                all_texts.extend([page.page_content for page in pages])
                logger.debug(f"Processed PDF: {file_path}")
            except Exception as e:
                logger.error(f"Error processing PDF {file_path}: {e}")
    
    # Create chunks and vectorize
    chunks = text_splitter.create_documents(all_texts)
    uuids = [str(uuid4()) for _ in range(len(chunks))]
    
    # Initialize and save vector store
    persist_dir = os.path.join(kb_dir, "vector_store")
    ensure_dir_exists(persist_dir)
    
    vectorstore = Chroma(
        collection_name=kb_name,
        embedding_function=embeddings,
        persist_directory=persist_dir,
    )
    
    # Add documents to the vector store
    try:
        vectorstore.add_documents(documents=chunks, ids=uuids)
        print(f"Added {len(chunks)} document chunks to vector store")
        logger.info(f"Added {len(chunks)} document chunks to vector store")
    except Exception as e:
        logger.error(f"Error adding documents to vector store: {e}")

def generate_generic_report(retriever, company_name, kb_name):
    """Generate a generic report for other companies"""
    logger.info(f"Generating generic report for {company_name}/{kb_name}")
    
    # Create the PowerPoint presentation
    ppt = create_presentation(company_name)
    
    # Run generic queries and fill the presentation
    try:
        # Company overview
        company_overview = run_query(
            f"Provide a comprehensive overview of {company_name} based on the documents.",
            retriever
        )
        add_slide(ppt, "Company Overview", company_overview)
        
        # Key findings
        key_findings = run_query(
            f"What are the key findings or insights about {company_name} from the documents?",
            retriever
        )
        add_slide(ppt, "Key Findings", key_findings)
        
        # Recommendations
        recommendations = run_query(
            f"Based on the documents, what recommendations would you provide for {company_name}?",
            retriever
        )
        add_slide(ppt, "Recommendations", recommendations)
        
        # Save the presentation
        save_presentation(ppt, company_name, kb_name)
        
    except Exception as e:
        logger.error(f"Error generating generic report: {e}")



def generate_chat_response(query, company_name, kb_name):
    """Generate a conversational response based on knowledge base content"""
    try:
        logger.info(f"Generating chat response for query: '{query}' using {company_name}/{kb_name}")

         # Check if KB name is None or empty
        if not kb_name:
            return "Please select a knowledge base before asking questions."
        
         # Convert kb_id to kb_name if an ID was passed
        if kb_name and kb_name.isdigit():
            try:
                from chatbot.models import KnowledgeBase
                kb_obj = KnowledgeBase.objects.get(id=int(kb_name))
                kb_name = kb_obj.name
                logger.info(f"Converted kb_id {kb_name} to name: {kb_obj.name}")
            except Exception as e:
                logger.error(f"Error converting kb_id to name: {e}")
                return "There was an error with the knowledge base selection. Please try again."
        
        # Check if KB name is None or empty
        if not kb_name:
            return "Please select a knowledge base before asking questions."
            
        # Check if company name is None or empty
        if not company_name:
            return "Please select a company before asking questions."
        
        # Check if knowledge base exists
        kb_dir = get_kb_dir(company_name, kb_name)
        persist_dir = os.path.join(kb_dir, "vector_store")
        
        if not os.path.exists(persist_dir):
            logger.warning(f"Vector store for {kb_name} does not exist")
            return "I don't have any information about this topic. Please upload relevant documents first."
        
        # Ensure LLM is initialized
        if not llm:
            logger.error("LLM not initialized for chat")
            return "I'm unable to process your request right now. Please try again later."
        
        # Load vector store and create retriever
        try:
            vectorstore = Chroma(
                collection_name=kb_name,
                embedding_function=embeddings,
                persist_directory=persist_dir,
            )
            
            retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 3})
        except Exception as e:
            logger.error(f"Error loading vector store: {e}")
            return "I'm having trouble accessing the knowledge base. Please try again later."
        
        # Get relevant documents with proper error handling
        try:
            retrieved_docs = retriever.invoke(query)
        except Exception as e:
            logger.error(f"Error retrieving documents: {e}")
            return "I encountered an error while searching for information. Please try again later."
        
        if not retrieved_docs:
            logger.warning(f"No relevant documents found for query: '{query}'")
            return "I couldn't find specific information to answer your question. Could you rephrase or ask about a different topic?"
        
        # Combine document content for context with safety check
        valid_contents = []
        for doc in retrieved_docs:
            if hasattr(doc, 'page_content') and doc.page_content is not None:
                valid_contents.append(doc.page_content)
        
        if not valid_contents:
            logger.warning("Retrieved documents had no valid content")
            return "I found some information but couldn't process it correctly. Please try a different question."
            
        context = "\n\n".join(valid_contents)
        
        # Create prompt template
        prompt = PromptTemplate(
            template="""You are a professional business analyst assistant.
            
            Based on the following context information, provide a well-structured, professional response to the query.
            Make your response comprehensive yet concise.
            Format your response with bullet points where appropriate.
            If you don't know the answer based on the context, say so clearly rather than making up information.
            
            Context: {context}
            
            Query: {query}
            
            Response:""",
            input_variables=["context", "query"]
        )
        
        # Create and run chain with timeout handling
        try:
            chain = LLMChain(llm=llm, prompt=prompt)
            response = chain.invoke({"context": context, "query": query})
            
            if not response or 'text' not in response:
                logger.error("Received empty or invalid response from LLM")
                return "I'm sorry, but I couldn't generate a proper response. Please try again later."
                
            return response['text'].strip()
        except Exception as chain_error:
            logger.error(f"Error in LLM chain: {chain_error}")
            return "I encountered an error while processing your request. Please try again later."
        
    except Exception as e:
        import traceback
        logger.error(f"Error generating chat response: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        return "I encountered an error while processing your request. Please try again later."

def run_query(query, retriever):
    """Run a query against the retriever and process with LLM"""
    if not llm:
        return "LLM not initialized. Cannot generate report content."
    
    # Get relevant documents
    retrieved_docs = retriever.invoke(query)
    context = "\n\n".join([doc.page_content for doc in retrieved_docs])
    
    # Create prompt template
    prompt = PromptTemplate(
        template="""You are a professional business analyst creating a report.
        
        Based on the following context information, provide a well-structured, professional response to the query.
        Make your response comprehensive yet concise, suitable for a business presentation slide.
        Format your response with bullet points where appropriate.
        
        Context: {context}
        
        Query: {query}
        
        Response:""",
        input_variables=["context", "query"]
    )
    
    # Create and run chain
    chain = LLMChain(llm=llm, prompt=prompt)
    response = chain.invoke({"context": context, "query": query})
    return response['text'].strip()

def create_presentation(company_name):
    """Create a new PowerPoint presentation"""
    try:
        # Try to use template if it exists
        if os.path.exists(PPT_TEMPLATE_PATH):
            ppt = Presentation(PPT_TEMPLATE_PATH)
        else:
            # Create a new presentation if template doesn't exist
            ppt = Presentation()
            
            # Add title slide
            title_slide_layout = ppt.slide_layouts[0]
            slide = ppt.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = f"{company_name} Analysis Report"
            subtitle.text = f"Generated on {datetime.datetime.now().strftime('%Y-%m-%d')}"
        
        return ppt
    except Exception as e:
        logger.error(f"Error creating presentation: {e}")
        return None

def add_slide(ppt, title, content):
    """Add a new slide to the presentation"""
    try:
        # Use bullet slide layout
        bullet_slide_layout = ppt.slide_layouts[1]
        slide = ppt.slides.add_slide(bullet_slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set content
        content_shape = slide.placeholders[1]
        content_shape.text = content
        
    except Exception as e:
        logger.error(f"Error adding slide: {e}")

def save_presentation(ppt, company_name, kb_name):
    """Save the presentation to file"""
    if not ppt:
        logger.error("Cannot save presentation: presentation object is None")
        return None
        
    try:
        # Create directory for the presentation
        ppt_dir = get_ppt_dir(company_name, kb_name)
        ppt_filename = f"{company_name}_{kb_name}_report.pptx"
        ppt_path = os.path.join(ppt_dir, ppt_filename)
        
        # Save the presentation
        ppt.save(ppt_path)
        logger.info(f"Presentation saved to {ppt_path}")
        
        pdf_path = None
        try:
            pdf_path = convert_ppt_to_pdf(ppt_path, company_name, kb_name)
        except Exception as e:
            logger.error(f"PDF conversion failed but presentation was saved: {e}")
        
        # Return the PowerPoint path regardless of PDF conversion success
        return ppt_path
    except Exception as e:
        logger.error(f"Error saving presentation: {e}")
        return None

# Update the convert_ppt_to_pdf function

def convert_ppt_to_pdf(ppt_path, company_name, kb_name):
    """Convert PowerPoint to PDF"""
    try:
        # Create directory for PDF output
        pdf_dir = get_pdf_dir(company_name, kb_name)
        pdf_filename = os.path.basename(ppt_path).replace('.pptx', '.pdf')
        pdf_path = os.path.join(pdf_dir, pdf_filename)
        
        # Try multiple conversion methods
        methods_tried = []
        
        # Method 1: Use LibreOffice if available (platform-independent)
        try:
            import subprocess
            methods_tried.append("LibreOffice")
            
            # Check for LibreOffice on different platforms
            libreoffice_commands = [
                'libreoffice', 'soffice',  # Linux/Mac
                r'C:\Program Files\LibreOffice\program\soffice.exe',  # Windows default
            ]
            
            cmd = None
            for command in libreoffice_commands:
                try:
                    subprocess.run([command, '--version'], stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False)
                    cmd = command
                    break
                except FileNotFoundError:
                    continue
                    
            if cmd:
                # Convert using LibreOffice
                result = subprocess.run([
                    cmd, '--headless', '--convert-to', 'pdf',
                    '--outdir', pdf_dir, ppt_path
                ], stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False)
                
                if result.returncode == 0:
                    logger.info(f"Successfully converted to PDF using LibreOffice: {pdf_path}")
                    return pdf_path
                else:
                    logger.warning(f"LibreOffice conversion failed: {result.stderr.decode()}")
            else:
                logger.warning("LibreOffice not found")
                
        except Exception as e:
            logger.warning(f"Failed to convert using LibreOffice: {e}")
        
        # Method 2: Using comtypes on Windows
        if os.name == 'nt':  # Only try on Windows
            try:
                methods_tried.append("comtypes")
                # Dynamic import to avoid errors on non-Windows platforms
                import importlib.util
                spec = importlib.util.find_spec('comtypes')
                
                if spec is not None:
                    import comtypes.client
                    import time
                    
                    # Get absolute paths
                    ppt_path_abs = os.path.abspath(ppt_path)
                    pdf_path_abs = os.path.abspath(pdf_path)
                    
                    # Initialize PowerPoint
                    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
                    powerpoint.Visible = False
                    
                    # Open and convert
                    presentation = powerpoint.Presentations.Open(ppt_path_abs)
                    presentation.SaveAs(pdf_path_abs, 32)  # 32 = PDF format
                    presentation.Close()
                    powerpoint.Quit()
                    
                    if os.path.exists(pdf_path):
                        logger.info(f"Successfully converted to PDF using comtypes: {pdf_path}")
                        return pdf_path
                else:
                    logger.warning("comtypes package not available")
            except Exception as e:
                logger.warning(f"Failed to convert using comtypes: {e}")
        
        # Method 3: Try using win32com as another option on Windows
        if os.name == 'nt':
            try:
                methods_tried.append("win32com")
                import win32com.client
                import pythoncom
                
                # Initialize COM in this thread
                pythoncom.CoInitialize()
                
                # Convert using PowerPoint via COM
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = False
                
                deck = powerpoint.Presentations.Open(os.path.abspath(ppt_path))
                deck.SaveAs(os.path.abspath(pdf_path), 32)
                deck.Close()
                powerpoint.Quit()
                
                # Clean up COM
                pythoncom.CoUninitialize()
                
                if os.path.exists(pdf_path):
                    logger.info(f"Successfully converted to PDF using win32com: {pdf_path}")
                    return pdf_path
            except Exception as e:
                logger.warning(f"Failed to convert using win32com: {e}")
        
        # If all previous methods failed, just copy the PowerPoint file to the output directory
        logger.warning(f"All conversion methods failed: {', '.join(methods_tried)}. Copying PPTX file instead.")
        shutil.copy2(ppt_path, os.path.join(pdf_dir, os.path.basename(ppt_path)))
        
        # Return the original PowerPoint path
        return ppt_path
        
    except Exception as e:
        logger.error(f"Error in PowerPoint to PDF conversion: {e}")
        return ppt_path  # Return original file path as fallback

def generate_report_from_files(company_name, kb_name, files):
    """Entry point function to generate report from uploaded files"""
    try:
        # Create vector store from uploaded files
        retriever = create_vector_store(files, kb_name, company_name)
        
        # Get paths to the generated reports
        ppt_dir = get_ppt_dir(company_name, kb_name)
        ppt_filename = f"{company_name}_{kb_name}_report.pptx"
        ppt_path = os.path.join(ppt_dir, ppt_filename)
        
        # Get PDF path
        pdf_dir = get_pdf_dir(company_name, kb_name)
        pdf_filename = ppt_filename.replace('.pptx', '.pdf')
        pdf_path = os.path.join(pdf_dir, pdf_filename)
        
        # Check if PDF was successfully created
        if os.path.exists(pdf_path):
            return {
                'success': True,
                'report_path': pdf_path,
                'message': f"Generated PDF report for {company_name}"
            }
        else:
            # Fall back to PPTX if PDF conversion failed
            return {
                'success': True,
                'report_path': ppt_path,
                'message': f"Generated PowerPoint report for {company_name} (PDF conversion failed)"
            }
    except Exception as e:
        logger.error(f"Error generating report from files: {e}")
        return {
            'success': False,
            'message': f"Error generating report: {str(e)}"
        }


# Add a new helper function to process files       
def process_files_for_report(company_name, kb_name, file_paths):
    """Process files and generate a report"""
    try:
        # Process PDF files
        all_texts = []
        for file_path in file_paths:
            if file_path.lower().endswith('.pdf'):
                try:
                    loader = PyPDFLoader(file_path)
                    pages = loader.load_and_split()
                    all_texts.extend([page.page_content for page in pages])
                    logger.debug(f"Processed PDF: {file_path}")
                except Exception as e:
                    logger.error(f"Error processing PDF {file_path}: {e}")
        
        # Create chunks and vectorize
        chunks = text_splitter.create_documents(all_texts)
        
        # Get directories
        kb_dir = get_kb_dir(company_name, kb_name)
        persist_dir = os.path.join(kb_dir, "vector_store")
        ensure_dir_exists(persist_dir)
        
        # Initialize vector store
        vectorstore = Chroma(
            collection_name=kb_name,
            embedding_function=embeddings,
            persist_directory=persist_dir,
        )
        
        # Add documents to the vector store in smaller batches
        # Use a batch size of 100 which is safer than the 166 limit
        batch_size = 100
        for i in range(0, len(chunks), batch_size):
            batch_end = min(i + batch_size, len(chunks))
            batch_chunks = chunks[i:batch_end]
            batch_uuids = [str(uuid4()) for _ in range(len(batch_chunks))]
            
            vectorstore.add_documents(documents=batch_chunks, ids=batch_uuids)
            logger.info(f"Added batch {i//batch_size + 1} with {len(batch_chunks)} document chunks to vector store")
        
        logger.info(f"Total added: {len(chunks)} document chunks to vector store")
        
        # Create retriever
        retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 3})
        
        # Generate report based on company
        if company_name == "Central Bank of Bahrain":
            generate_cbb_report(retriever, company_name, kb_name)
        else:
            generate_generic_report(retriever, company_name, kb_name)
        
        # Return the path to the generated report
        ppt_dir = get_ppt_dir(company_name, kb_name)
        ppt_filename = f"{company_name}_{kb_name}_report.pptx"
        ppt_path = os.path.join(ppt_dir, ppt_filename)
        
        return {
            'success': True,
            'report_path': ppt_path,
            'message': f"Generated report for {company_name}"
        }
    except Exception as e:
        logger.error(f"Error processing files for report: {e}")
        return {
            'success': False,
            'message': f"Error processing files: {str(e)}"
        }


def generate_report_from_kb(company_name, kb_name):
    """Entry point function to generate report from existing knowledge base"""
    try:
        # Get the vector store directory path
        kb_dir = get_kb_dir(company_name, kb_name)
        persist_dir = os.path.join(kb_dir, "vector_store")
        
        if not os.path.exists(persist_dir):
            logger.warning(f"Vector store for {kb_name} does not exist. Attempting to create it.")
            
            # Check if we have files to process for this KB
            from .models import KnowledgeBase, KnowledgeBaseFile
            try:
                kb = KnowledgeBase.objects.get(name=kb_name, company=company_name)
                kb_files = KnowledgeBaseFile.objects.filter(knowledge_base=kb)
                
                if kb_files.exists():
                    # Process files to create vector store
                    # Fix: Properly handle file access by saving them to a temporary location first
                    temp_files = []
                    for file_obj in kb_files:
                        # Create temp directory if it doesn't exist
                        temp_dir = os.path.join(kb_dir, "temp")
                        ensure_dir_exists(temp_dir)
                        
                        # Get original filename
                        original_filename = os.path.basename(file_obj.file.name)
                        
                        # Create a temp file path
                        temp_file_path = os.path.join(temp_dir, original_filename)
                        
                        # Save the file content to the temp location
                        with open(temp_file_path, 'wb') as temp_file:
                            for chunk in file_obj.file.chunks():
                                temp_file.write(chunk)
                        
                        temp_files.append(temp_file_path)
                    
                    # Now process these temporary files
                    return process_files_for_report(company_name, kb_name, temp_files)
                else:
                    return {
                        'success': False,
                        'message': f"No files found for knowledge base {kb_name}. Please upload files first."
                    }
            except Exception as inner_e:
                logger.error(f"Error accessing knowledge base files: {str(inner_e)}")
                return {
                    'success': False,
                    'message': f"Vector store for {kb_name} does not exist and could not create it: {str(inner_e)}"
                }
        
        # Load the vector store
        vectorstore = Chroma(
            collection_name=kb_name,
            embedding_function=embeddings,
            persist_directory=persist_dir,
        )
        
        # Create retriever
        retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 3})
        
        # Generate report based on company
        if company_name == "Central Bank of Bahrain":
            generate_cbb_report(retriever, company_name, kb_name)
        else:
            generate_generic_report(retriever, company_name, kb_name)
        
        # Return the path to the generated report
        ppt_dir = get_ppt_dir(company_name, kb_name)
        ppt_filename = f"{company_name}_{kb_name}_report.pptx"
        ppt_path = os.path.join(ppt_dir, ppt_filename)
        
        return {
            'success': True,
            'report_path': ppt_path,
            'message': f"Generated report for {company_name} from knowledge base {kb_name}"
        }
        
    except Exception as e:
        logger.error(f"Error generating report from KB: {e}")
        return {
            'success': False,
            'message': f"Error generating report: {str(e)}"
        }
    

